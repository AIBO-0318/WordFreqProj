# -*- coding: utf-8 -*-
"""왓챠 코멘트 분석 AI 서비스 — 발표자료(10장) 생성 (python-pptx)"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 팔레트 ────────────────────────────────────────────────
INK    = "1B2333"   # 다크 네이비 (표지/마무리 배경, 진한 텍스트)
ACCENT = "FF3D71"   # 왓챠풍 코랄핑크 (강조)
POS    = "2ECC71"   # 긍정 초록
NEG    = "E74C3C"   # 부정 빨강
CARD   = "F3F5F9"   # 카드 배경(라이트)
CARD2  = "FFFFFF"
MUTED  = "6B7280"
TEXT   = "222B3A"
WHITE  = "FFFFFF"
LINE   = "E2E7EF"
FONT   = "Malgun Gothic"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def hex2rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = hex2rgb(color)


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, rounded=False, shadow=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = hex2rgb(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = hex2rgb(line); shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    if shadow:
        el = shp._element.spPr
        sp = el.makeelement(qn('a:effectLst'), {})
        sh = sp.makeelement(qn('a:outerShdw'),
                            {'blurRad': '60000', 'dist': '25000', 'dir': '5400000', 'rotWithShape': '0'})
        clr = sh.makeelement(qn('a:srgbClr'), {'val': '9AA3B2'})
        al = clr.makeelement(qn('a:alpha'), {'val': '40000'})
        clr.append(al); sh.append(clr); sp.append(sh); el.append(sp)
    return shp


def txt(s, runs, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each para = list of (text, size, color, bold) tuples."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = hex2rgb(col); r.font.name = FONT
            rPr = r._r.get_or_add_rPr()
            rFonts = rPr.find(qn('a:rFonts'))
            if rFonts is None:
                rFonts = rPr.makeelement(qn('a:rFonts'), {})
                rPr.insert(0, rFonts)
            rFonts.set(qn('a:eastAsia'), FONT)
            rFonts.set(qn('a:latin'), FONT)
    return tb


def title_block(s, num, title, sub=None):
    """라이트 콘텐츠 슬라이드 공통 제목 (왼쪽 코랄 사각 모티프 + 번호)."""
    rect(s, 0.6, 0.62, 0.13, 0.62, fill=ACCENT)
    txt(s, [[(title, 30, INK, True)]], 0.85, 0.55, 11.5, 0.8, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        txt(s, [[(sub, 14, MUTED, False)]], 0.86, 1.28, 11.5, 0.4)
    txt(s, [[(f"{num:02d}", 12, MUTED, True)]], SW - 1.2, 6.95, 0.6, 0.4, align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, head, head_color, lines, fill=CARD):
    rect(s, x, y, w, h, fill=fill, rounded=True, shadow=True)
    pad = 0.28
    runs = [[(head, 15, head_color, True)]]
    for ln in lines:
        runs.append([(ln, 12.5, TEXT, False)])
    txt(s, runs, x + pad, y + 0.22, w - 2 * pad, h - 0.4, space_after=5, line_spacing=1.05)


# ════════════════════════════════════════════════════════
# 1. 표지 (다크)
# ════════════════════════════════════════════════════════
s = slide(); bg(s, INK)
rect(s, 0, 0, 0.28, SH, fill=ACCENT)
txt(s, [[("AI 서비스 개발 프로젝트", 16, ACCENT, True)]], 1.1, 1.7, 10, 0.5)
txt(s, [
    [("왓챠피디아 코멘트 기반", 40, WHITE, True)],
    [("영화 여론 분석 AI 서비스", 40, WHITE, True)],
], 1.1, 2.4, 11, 1.8, line_spacing=1.05)
txt(s, [[("단어 빈도 분석  ·  감성 분석(LSTM)  ·  영화 추천(TF-IDF)", 17, "C9D2E0", False)]],
    1.12, 4.35, 11, 0.5)
txt(s, [[("인공지능서비스개발   |   [팀명 / 발표자]   |   2026", 13, MUTED, False)]],
    1.12, 6.5, 11, 0.4)

# ════════════════════════════════════════════════════════
# 2. 프로젝트 개요 & 니즈
# ════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
title_block(s, 2, "프로젝트 개요 & 니즈", "왜 만들었나")
txt(s, [
    [("영화 리뷰는 넘쳐나지만 흩어져 있다", 18, INK, True)],
    [("한 영화에 수천~수만 개의 코멘트 → 사람이 다 읽고 여론을 파악하기 어렵다.", 14, TEXT, False)],
], 0.85, 1.9, 6.2, 1.3, space_after=8)
txt(s, [
    [("→ 코멘트를 ", 15, TEXT, False), ("자동으로 분석·시각화", 15, ACCENT, True),
     ("하여", 15, TEXT, False)],
    [("   ① 무슨 단어가 많은지  ② 호불호(감성)  ③ 비슷한 영화", 14, TEXT, False)],
    [("   를 한눈에 보여주는 웹 대시보드를 만든다.", 14, TEXT, False)],
], 0.85, 3.3, 6.2, 1.8, space_after=6)

# 우측 stat 카드
sx = 7.6
for i, (val, lab) in enumerate([("150만+", "수집한 코멘트"), ("309편", "영화"), ("1957–2025", "개봉연도 범위")]):
    cy = 1.95 + i * 1.5
    rect(s, sx, cy, 5.1, 1.3, fill=CARD, rounded=True, shadow=True)
    txt(s, [[(val, 34, ACCENT, True)]], sx + 0.35, cy + 0.18, 3.0, 0.95, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, [[(lab, 14, MUTED, True)]], sx + 3.0, cy + 0.18, 1.9, 0.95, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
txt(s, [[("데이터 출처: 왓챠피디아 직접 크롤링", 11, MUTED, False)]], sx, 6.4, 5.1, 0.4, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════
# 3. 서비스 구성
# ════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
title_block(s, 3, "서비스 구성", "데이터 → 분석 모듈 → 웹 대시보드")
# 파이프라인
flow = ["왓챠 크롤링\n(코멘트 150만)", "전처리·형태소\n(Okt 토큰화)", "분석 모듈\n(mylib 패키지)", "Streamlit\n대시보드"]
fw, gap = 2.7, 0.45
fx = 0.85
for i, step in enumerate(flow):
    x = fx + i * (fw + gap)
    rect(s, x, 1.95, fw, 1.1, fill=INK if i in (0, 3) else CARD, rounded=True, shadow=True)
    head, sub = step.split("\n")
    col = WHITE if i in (0, 3) else INK
    subcol = "C9D2E0" if i in (0, 3) else MUTED
    txt(s, [[(head, 14, col, True)], [(sub, 11, subcol, False)]],
        x + 0.15, 2.05, fw - 0.3, 0.95, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=2)
    if i < 3:
        txt(s, [[("›", 22, ACCENT, True)]], x + fw + 0.02, 2.05, gap, 0.95,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 3개 기능 카드
feats = [("📊  단어 빈도 분석", ACCENT, ["코멘트에서 명사 추출", "워드클라우드 · 막대그래프", "영화/연도/전체 단위"]),
         ("😊  감성 분석", POS, ["긍정/부정 자동 분류", "LSTM 딥러닝 모델", "긍·부정 단어 따로 시각화"]),
         ("🎬  영화 추천", "3B82F6", ["코멘트 유사도 기반", "TF-IDF + 코사인", "비슷한 영화 Top N"])]
cw = 3.95
for i, (h, c, lines) in enumerate(feats):
    x = 0.85 + i * (cw + 0.3)
    card(s, x, 3.65, cw, 2.7, h, c, lines)

# ════════════════════════════════════════════════════════
# 4. ML/DL 적용 부분
# ════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
title_block(s, 4, "어디에 머신러닝·딥러닝을 적용했나")
rows = [
    ("딥러닝 (DL)", NEG_:=ACCENT, "감성 분석",
     "순환신경망 LSTM — Embedding → LSTM(64) → Dense → Softmax. 코멘트를 긍정/부정으로 분류."),
    ("머신러닝 (ML)", "3B82F6", "영화 추천",
     "TF-IDF로 영화별 코멘트를 벡터화하고 코사인 유사도로 비슷한 영화를 계산."),
    ("자연어처리 (NLP)", POS, "형태소 분석",
     "KoNLPy Okt로 한국어를 형태소 단위로 토큰화 · 명사 추출 (모든 기능의 기반)."),
]
for i, (tag, c, name, desc) in enumerate(rows):
    y = 1.95 + i * 1.35
    rect(s, 0.85, y, 12.0, 1.15, fill=CARD, rounded=True, shadow=True)
    rect(s, 0.85, y, 0.13, 1.15, fill=c)
    txt(s, [[(tag, 16, c, True)], [(name, 12, MUTED, True)]], 1.15, y + 0.18, 2.7, 0.8,
        anchor=MSO_ANCHOR.MIDDLE, space_after=2)
    txt(s, [[(desc, 13.5, TEXT, False)]], 4.0, y + 0.18, 8.6, 0.8, anchor=MSO_ANCHOR.MIDDLE)
txt(s, [[("※ 단어 빈도 분석은 통계 기반 (빈도 집계)으로, 별도 학습 없이 동작", 12, MUTED, False)]],
    0.85, 6.2, 12, 0.4)

# ════════════════════════════════════════════════════════
# 5. 데이터 준비 & 처리
# ════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
title_block(s, 5, "데이터 준비 & 처리")
steps = [("크롤링", "왓챠 API로 코멘트\n150만 건 수집"),
         ("전처리", "중복·결측 제거,\n정규화·정제"),
         ("자동 라벨링", "별점이 없어\n감성사전으로 라벨 생성"),
         ("학습 데이터", "긍·부정 균형\n약 12.8만 건")]
fw, gap = 2.85, 0.35
for i, (h, d) in enumerate(steps):
    x = 0.85 + i * (fw + gap)
    rect(s, x, 2.0, fw, 1.5, fill=CARD, rounded=True, shadow=True)
    rect(s, x, 2.0, fw, 0.5, fill=INK, rounded=True)
    txt(s, [[(h, 15, WHITE, True)]], x, 2.05, fw, 0.45, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, [[(d.split(chr(10))[0], 13, TEXT, True)], [(d.split(chr(10))[1], 12, MUTED, False)]],
        x + 0.2, 2.6, fw - 0.4, 0.85, align=PP_ALIGN.CENTER, space_after=2)
    if i < 3:
        txt(s, [[("›", 20, ACCENT, True)]], x + fw, 2.0, gap, 1.5, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, 0.85, 4.15, 12.0, 1.9, fill="FFF0F4", rounded=True)
txt(s, [
    [("핵심 아이디어 — 별점이 없는데 어떻게 정답(라벨)을 만드나?", 16, ACCENT, True)],
    [("왓챠 코멘트에는 별점이 저장돼 있지 않다. 그래서 ", 14, TEXT, False),
     ("KNU 한국어 감성사전", 14, INK, True),
     ("으로 코멘트의", 14, TEXT, False)],
    [("긍정/부정을 자동 판정해 학습용 라벨을 생성했다. (네이버 NSMC가 별점으로 라벨을 만든 것과 같은 발상)", 14, TEXT, False)],
], 1.2, 4.4, 11.3, 1.5, space_after=6, line_spacing=1.05)

# ════════════════════════════════════════════════════════
# 6. 문제 해결 ① — 라벨 품질 끌어올리기
# ════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
title_block(s, 6, "문제 해결 과정 ① — 감성 라벨 만들기", "자동 라벨링의 오류를 단계적으로 교정")
items = [
    ("정답(별점)이 없음", "KNU 감성사전으로 긍·부정 자동 라벨링"),
    ("조사 '이/가'를 긍정 단어로 오인 → 편향", "내용어(명사·동사·형용사)만 점수화"),
    ("\"재미없다\" 같은 부정 표현을 무시", "부정어 처리 — 앞 단어의 극성을 반전"),
    ("'재밌다·꿀잼' 등 구어체를 사전이 모름", "영화 리뷰 구어체 감성어를 직접 보강"),
]
for i, (prob, sol) in enumerate(items):
    y = 2.0 + i * 1.08
    rect(s, 0.85, y, 5.85, 0.92, fill="FDECEC", rounded=True)
    txt(s, [[("문제", 11, NEG, True)], [(prob, 13, TEXT, False)]], 1.1, y + 0.13, 5.4, 0.7, space_after=2)
    txt(s, [[("→", 20, INK, True)]], 6.75, y, 0.5, 0.92, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 7.3, y, 5.55, 0.92, fill="EAF7EF", rounded=True)
    txt(s, [[("해결", 11, POS, True)], [(sol, 13, TEXT, False)]], 7.55, y + 0.13, 5.1, 0.7, space_after=2)

# ════════════════════════════════════════════════════════
# 7. 문제 해결 ② — 편향과 진단
# ════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
title_block(s, 7, "문제 해결 과정 ② — 편향 교정과 진단", "\"여기가 문제라 데이터를 더 모아봤다\"")
card(s, 0.85, 1.95, 5.85, 1.9, "주제어·분위기어로 호평이 부정 처리", NEG,
     ["'가난·죽음·무섭다'가 들어간 호평이", "부정으로 오분류 →", "주제어·분위기어를 점수에서 제외 + 임계값 보정"])
card(s, 7.0, 1.95, 5.85, 1.9, "데이터를 7만 → 14만으로 2배 증량", "3B82F6",
     ["\"데이터가 부족한가?\" 가설 검증을 위해", "샘플 30만→60만으로 재학습 →", "정확도 거의 변화 없음"])
rect(s, 0.85, 4.15, 12.0, 1.95, fill=INK, rounded=True, shadow=True)
txt(s, [[("진단", 13, ACCENT, True)]], 1.2, 4.4, 11, 0.4)
txt(s, [
    [("병목은 데이터의 '양'이 아니라 ", 22, WHITE, True), ("라벨의 '품질'", 22, ACCENT, True)],
    [("감성사전은 '슬프다·무섭다'가 칭찬인지 비판인지 구분하지 못한다. → 근본 해결책은 코멘트별 실제 별점 라벨.", 14, "C9D2E0", False)],
], 1.2, 4.85, 11.3, 1.1, space_after=8, line_spacing=1.05)

# ════════════════════════════════════════════════════════
# 8. 성능 평가
# ════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
title_block(s, 8, "성능 평가")
# 좌: 감성
rect(s, 0.85, 1.95, 5.85, 4.3, fill=CARD, rounded=True, shadow=True)
txt(s, [[("😊  감성 분석 (LSTM)", 17, POS, True)]], 1.15, 2.2, 5.3, 0.5)
txt(s, [
    [("• 테스트 정확도 ", 13.5, TEXT, False), ("약 90%", 13.5, INK, True), (" (자동 라벨 기준)", 13.5, MUTED, False)],
    [("• 명확한 문장은 정확", 13.5, TEXT, False)],
    [("   \"시간이 아까움\" → 부정 / \"최고의 명작\" → 긍정", 12.5, MUTED, False)],
    [("• 한계: 진지한 영화의 호평 일부를 부정으로 오분류", 13.5, TEXT, False)],
    [("• 확신도 ≠ 정확도 (과신 경향) → 지표에서 제외", 13.5, TEXT, False)],
], 1.15, 2.8, 5.4, 3.2, space_after=10, line_spacing=1.05)
# 우: 추천 (정성)
rect(s, 7.0, 1.95, 5.85, 4.3, fill=CARD, rounded=True, shadow=True)
txt(s, [[("🎬  영화 추천 (정성 평가)", 17, "3B82F6", True)]], 7.3, 2.2, 5.3, 0.5)
rec = [("기생충", "설국열차·미키17·마더 (봉준호)"),
       ("인터스텔라", "테넷·인셉션·오펜하이머 (놀란)"),
       ("어벤져스:엔드게임", "마블 시리즈 (유사도 0.70)")]
for i, (m, r) in enumerate(rec):
    y = 2.95 + i * 1.0
    txt(s, [[(m, 14, INK, True)], [(r, 12.5, TEXT, False)]], 7.3, y, 5.3, 0.9, space_after=2)
txt(s, [[("코멘트 단어만으로 감독·시리즈를 정확히 포착", 12, MUTED, False)]], 7.3, 5.95, 5.3, 0.3)

# ════════════════════════════════════════════════════════
# 9. 시연
# ════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
title_block(s, 9, "시연", "localhost:8501 — 실시간 동작")
shots = ["📊 단어 빈도 분석 탭", "😊 감성 분석 탭", "🎬 영화 추천 탭"]
cw = 3.95
for i, cap in enumerate(shots):
    x = 0.85 + i * (cw + 0.3)
    rect(s, x, 2.0, cw, 3.3, fill=CARD, line=LINE, lw=1, rounded=True)
    txt(s, [[("[ 스크린샷 ]", 13, MUTED, True)]], x, 3.3, cw, 0.5, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, [[(cap, 14, INK, True)]], x, 5.45, cw, 0.5, align=PP_ALIGN.CENTER)
txt(s, [[("※ 발표 시 캡처 이미지를 위 영역에 삽입", 11, MUTED, False)]], 0.85, 6.4, 12, 0.4)

# ════════════════════════════════════════════════════════
# 10. 결론 & 한계·향후 (다크)
# ════════════════════════════════════════════════════════
s = slide(); bg(s, INK)
rect(s, 0.6, 0.75, 0.13, 0.62, fill=ACCENT)
txt(s, [[("결론 & 향후 과제", 30, WHITE, True)]], 0.85, 0.68, 11, 0.8, anchor=MSO_ANCHOR.MIDDLE)
pts = [
    ("3개의 AI 서비스 완성", "단어 빈도 · 감성 분석(LSTM) · 영화 추천(TF-IDF)을 하나의 대시보드로"),
    ("문제를 단계적으로 진단·해결", "라벨링 오류를 사전·내용어·부정어·구어체 처리로 차례로 교정"),
    ("한계를 실험으로 규명", "데이터 2배 실험 → 병목은 '양'이 아니라 '라벨 품질'임을 입증"),
    ("향후: 실제 별점 라벨", "코멘트별 별점으로 라벨링하면 정확도를 근본적으로 끌어올릴 수 있음"),
]
for i, (h, d) in enumerate(pts):
    y = 1.85 + i * 1.15
    rect(s, 0.85, y + 0.05, 0.5, 0.5, fill=ACCENT, rounded=True)
    txt(s, [[(str(i + 1), 18, WHITE, True)]], 0.85, y + 0.05, 0.5, 0.5, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, [[(h, 17, WHITE, True)], [(d, 13, "C9D2E0", False)]], 1.6, y, 11, 1.0, space_after=3)
txt(s, [[("개발: Claude Code(AI)와 협업하여 데이터 처리·모델 학습·대시보드 구현", 12, MUTED, False)]],
    0.85, 6.85, 12, 0.4)

OUT = "watcha_presentation.pptx"
prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
